#./mdtree/tree2ppt.py
import os
from io import BytesIO
import re
import markdown
from PIL.ImageQt import rgb
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.text.text import Font
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from .parser import parse_string, Out, Heading
from .utils import get_random_theme, get_random_file, get_title_file, get_dominant_color
from APC.config_parser import get_args_parser
from .agent_color import ColorAgent
from .agent_componets import ComponentAgent
from .agent_componets_half import ComponentAgentHalf
from .template import template_1,template_1_prompt
from .template_contents import create_styled_toc_slide
from .template_gen import template_generation
from .template_half import template_half, template_half_prompt
from .visual_components import DEFAULT_THEME_COLORS, add_main_text_box, add_styled_quote_box, add_key_data_callout, \
    add_decorative_divider, add_icon_title_bullets_component, add_triple_card_component, add_parallel_boxes_component
import time
import logging
import random
import requests # 用于下载图片
from PIL import Image # 用于获取图片尺寸

os.environ["OPENBLAS_NUM_THREADS"] = "4"
os.environ["OMP_NUM_THREADS"] = "4"
os.environ["MKL_NUM_THREADS"] = "4"
logging.basicConfig(
    filename='slide_data.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S',
    filemode='a' # 追加模式，如果文件不存在则创建
)


class Tree2PPT:
    prs: Presentation = None
    md_str: str = None
    out: Out = None
    tree: Heading = None
    theme: str = None

    def __init__(self, md_str):
        self.config = get_args_parser()
        start_time = time.time()
        self.init_pptx(self.config.iamge)
        output_dir = self.config.output_dir

        # 丰富内容
        # self.content_agent = CotentAgent()
        # md_str1 = self.content_agent(md_str)
        md_str1 = md_str

        self.init_markdown(md_str1)
        original_text = self.tree.text
        self.valid_text = self.sanitize_filename(original_text)

        # 控制颜色
        self.color_agent = ColorAgent()
        COVER_dominant_color = get_dominant_color(get_title_file(self.theme))
        CONTENT_dominant_color = get_dominant_color(get_random_file(self.theme))
        self.cover_theme_colors = self.color_agent(COVER_dominant_color, True)
        self.content_theme_colors = self.color_agent(CONTENT_dominant_color, False)  # content必须后于cover
        self.infra_color = self.color_agent(CONTENT_dominant_color, False, True, self.valid_text)  # INFRA_COLOR必须后于content

        self.traverse_tree(self.tree)

        path = output_dir + self.valid_text + '.pptx'
        self.prs.save(path)
        self.path = path

        end_time = time.time()

        # 计算并打印运行时间
        execution_time = end_time - start_time
        print(f"代码片段的运行时间为: {execution_time:.4f} 秒")
        pass

    def sanitize_filename(self,text):
        """
        替换掉文件名中不合法的字符为下划线 '_'
        """
        invalid_chars_pattern = r'[\\/:*?"<>|\x00-\x1f：]'  # 添加了全角冒号
        sanitized_text = re.sub(invalid_chars_pattern, '_', text)
        sanitized_text = re.sub(r'_+', '_', sanitized_text)  # 多个下划线替换为单个
        sanitized_text = sanitized_text.strip('_')  # 移除开头和结尾的下划线
        if not sanitized_text:  # 处理完全由非法字符组成的极端情况
            sanitized_text = "default_presentation_title"  # 默认文件名
        return sanitized_text

    def get_path(self):
        return self.path

    def init_pptx(self, theme=None):
        prs = Presentation()
        self.theme = get_random_theme(theme)  # 随机选择一个背景主题
        self.prs = prs  # 初始化一个新的、空的 PowerPoint 演示文稿对象

    def init_markdown(self, md_str):
        self.md_str = md_str
        self.out = parse_string(md_str)
        self.tree = self.out.main

    def traverse_tree(self, heading):
        """
        递归遍历 Markdown 树，为每个节点创建幻灯片。
        """
        section_start_level = 2 # 默认的章节开始级别
        if heading is not None:
            content = "" # 初始化内容字符串
            if heading.source is None or heading.source == '': # 如果节点本身没有源文本
                 if heading.children: # 检查是否有子节点
                    for child in heading.children: # 拼接子节点的文本作为内容
                         content += child.text + "\n"
            else:
                content = heading.source # 使用节点自身的源文本作为内容

            is_cover_type_node = (len(self.prs.slides) == 0) or \
                                 (heading.level is not None and heading.level == section_start_level)

            current_theme_colors = self.cover_theme_colors if is_cover_type_node else self.content_theme_colors

            # --- 创建 MD2Slide 实例时传入继承的颜色 ---
            slide_obj = MD2Slide(
                self.prs, self.theme, heading.text, content=content, tree=self.tree,
                heading_level=heading.level,
                section_start_level=section_start_level,
                font_color=current_theme_colors,  # 传入封面字体颜色
                infra_color=self.infra_color,
                config=self.config
            )
        else:
            return

        # --- 5. 递归调用时，始终传递相同的固定颜色 ---
        if heading.children:
            for child in heading.children:
                self.traverse_tree(child)

    def save_stream(self):
        """
        将演示文稿保存到内存流中。
        """
        stream = BytesIO()
        self.prs.save(stream)
        stream.seek(0) # 将流的位置重置到开头
        return stream


class MarkdownCategory:
    TITLE = "#"
    CONTENT = "<p>"
    picture = "!["
    pass


class MD2Slide:
    """
    这个类负责在演示文稿中创建单张幻灯片。
    """
    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "微软雅黑"
    font_title_size: Pt = Pt(26)
    font_content_size: Pt = Pt(20)

    font_title_color: rgb = None # 会在 __init__ 中设置
    font_content_color: rgb = None # 会在 __init__ 中设置
    infra_color:rgb = None

    heading_level: int = None  # Add attribute for heading level
    section_start_level: int = 2 # Default section start level
    background_image_path: str = None # 新增：用于保存背景图路径

    def __init__(self, presentation, theme_path, title, content,tree=None, *args, heading_level=None, section_start_level=2, font_color=None, infra_color=None, config=None,**kwargs):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[8])  # 添加一张新的幻灯片。布局 8 通常是空白布局。

        self.title = title
        self.content = content
        self.theme = theme_path
        self.heading_level = heading_level # Store the heading level
        self.section_start_level = section_start_level # Store the section start level criterion
        self.tree = tree  # 接收tree参数

        self.font_title_color = font_color[0]
        self.font_content_color = font_color[1]
        self.infra_color = infra_color

        self.config = config
        if self.config.template is True:
            self.template = template_1
            self.template_half = template_half
        else:
            self.template = None

        self.init_slide()
        self.init_font(**kwargs)  # 设置字体相关的默认值（字体名称、标题/内容字号、颜色）
        self.init_title()
        self.init_content()


    def init_slide(self):
        """
        初始化幻灯片背景。移除默认占位符，并根据幻灯片类型添加背景图片。
        """
        for placeholder in self.slide.placeholders:
            placeholder.element.getparent().remove(placeholder.element)  # 首先移除新幻灯片上所有默认的占位符（文本框等）。

        use_special_background = False
        slide_idx = len(self.presentation.slides) # 获取当前幻灯片索引（从1开始）
        # 条件1：是封面页（第一张幻灯片）
        # 条件2：是主要小节的起始页（标题级别符合要求）
        if slide_idx == 1 or \
           (self.heading_level is not None and self.heading_level == self.section_start_level):
            use_special_background = True

        try:
            if use_special_background: # 如果使用特殊背景
                bg_path = get_title_file(self.theme) # 获取章节标题背景图
            else: # 否则使用随机主题背景图
                bg_path = get_random_file(self.theme)

            if bg_path and os.path.exists(bg_path): # 如果路径有效且文件存在
                self.slide.shapes.add_picture( # 添加图片作为背景
                    bg_path, Inches(0), Inches(0), # 从左上角开始
                    width=self.presentation.slide_width, # 宽度为幻灯片宽度
                    height=self.presentation.slide_height # 高度为幻灯片高度
                )
                self.background_image_path = bg_path # 保存背景图片路径
            else:
                print(f"警告：背景图片未在 '{bg_path}' 找到。")
        except Exception as e:
            print(f"从主题 '{self.theme}' 添加背景图片时出错：{e}")

        # 确保背景图片被置于底层（如果其他方法未处理）
        if self.slide.shapes: # 如果幻灯片上有形状
             try:
                pass # 通常，首先将图片添加到空白幻灯片布局即可处理此问题。
             except Exception as e:
                print(f"无法将背景发送到后面：{e}")

    def init_font(self, **kwargs):
        """根据传入的参数设置字体名称和大小。"""
        if 'font_name' in kwargs: self.font_name = kwargs['font_name']
        if 'font_title_size' in kwargs: self.font_title_size = kwargs['font_title_size']
        if 'font_content_size' in kwargs: self.font_content_size = kwargs['font_content_size']

    # 根据是标题 (MarkdownCategory.TITLE) 还是内容 (MarkdownCategory.CONTENT)，
    # 为给定的 Font 对象设置相应的样式（字体、大小、颜色、加粗）。
    def get_font(self, font: Font, category: str):
        font.bold = True
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_toc_slide(self):
        """初始化目录页
            特别地将幻灯片格式化为目录（TOC）。
            添加标题 "目录"，然后遍历主 tree 根节点的子节点 (self.tree.children) 来创建列出主要章节的项目符号列表。
        """
        shapes = self.slide.shapes
        # 添加目录内容
        # 这里需要从tree中获取二级标题
        if hasattr(self, 'tree'):
            content = ""
            for child in self.tree.children:
                content += f"• {child.text}\n"

            create_styled_toc_slide(self.slide, content, self.infra_color)
            # content_box = shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
            # tf = content_box.text_frame
            # tf.clear()
            # paragraph = tf.paragraphs[0]
            # paragraph.text = content
            # self.get_font(paragraph.font, MarkdownCategory.CONTENT)

    def _process_title_text(self, title):
        """处理标题文本，在特定符号后添加换行符，以改善长标题的显示效果。"""
        break_chars = [':', '：', '-', '—'] # 定义需要换行的符号
        for char in break_chars:
            title = title.replace(char, char + '\n') # 替换
        return title

    def init_title(self):
        """
        向幻灯片添加标题文本框。
        根据当前幻灯片是第几张来决定标题的格式和位置：
            第一张（封面）: 标题居中，字号较大（Pt(48)），应用 _process_title_text 的换行处理。
            第二张（目录）: 标题固定为 "目录"，居中，使用标准标题字体。
            其他（内容页）: 标题位于幻灯片顶部，使用标准标题字体。
        :return:
        """
        shapes = self.slide.shapes
        # 封面页特殊处理
        if len(self.presentation.slides) == 1:  # 封面页
            text_box = self.slide.shapes.add_textbox(
                Inches(1),
                self.presentation.slide_height / 2 - Inches(1),
                self.presentation.slide_width - Inches(2),
                Inches(2))
            tf = text_box.text_frame
            tf.clear()
            tf.word_wrap = True
            title_text = self._process_title_text(self.title)
            paragraph = tf.paragraphs[0]
            paragraph.text = title_text
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            font = paragraph.font
            font.name = self.font_name
            font.size = Pt(48)  # 大号字体
            font.bold = True
            font.color.rgb = self.font_title_color
        elif len(self.presentation.slides) == 2:  # 目录页
            a=1
            # text_box = self.slide.shapes.add_textbox(
            #     Inches(1), Inches(0.8),
            #     self.presentation.slide_width - Inches(2), Inches(1))
            # tf = text_box.text_frame
            # tf.clear()
            # paragraph = tf.paragraphs[0]
            # paragraph.text = "目录"
            # paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            # self.get_font(paragraph.font, MarkdownCategory.TITLE)
        else:  # 内容页
            is_section_cover = self.heading_level is not None and self.heading_level == self.section_start_level
            if is_section_cover:  # 标签页
                # --- 章节封面标题布局：居中 ---
                # 可以调整 top 值，让标题稍微靠上一些
                left = Inches(0.5)  # 左右留边距
                top = Inches(0.7)   # 距离顶部的位置
                width = self.presentation.slide_width - Inches(1) # 宽度接近充满页面
                height = Inches(1.0) # 可以适当增加高度以防标题过长换行
                text_box = self.slide.shapes.add_textbox(left, top, width, height)
                tf = text_box.text_frame
                tf.clear()
                tf.word_wrap = True
                paragraph = tf.paragraphs[0]
                paragraph.text = self.title
                self.get_font(paragraph.font, MarkdownCategory.TITLE)
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER # <--- 水平居中
                # 可以考虑给章节封面标题设置稍大一点的字号
                # paragraph.font.size = Pt(32)
            else:  # 普通内容页
                left = Inches(0.3)
                top = Inches(0.7)
                width = Inches(9)
                height = Inches(0.8)
                text_box = self.slide.shapes.add_textbox(left, top, width, height)
                tf = text_box.text_frame
                tf.clear()
                tf.word_wrap = True
                paragraph = tf.paragraphs[0]
                paragraph.text = self.title
                self.get_font(paragraph.font, MarkdownCategory.TITLE)
                # 默认即为左对齐，无需显式设置 paragraph.alignment

    def init_content(self):
        """
        向幻灯片添加主要内容
        :return:
        """
        shapes = self.slide.shapes
        # 如果是第一页（封面），不需要内容，内容就是title，已经在init_title()方法里面初始化了
        if len(self.presentation.slides) == 1:
            return

        # 如果是第二页，创建目录页
        if len(self.presentation.slides) == 2:
            self.init_toc_slide()
            return

        is_section_cover = self.heading_level is not None and self.heading_level == self.section_start_level
        content_text = self.content.replace("<p>", "").replace("</p>", "\n")
        has_pic = False
        content_text, has_pic = self._add_images_from_content(content_text)  # 判断是否有图片

        if is_section_cover:
            # 章节封面布局：通常希望文本区域更居中，占满宽度
            left = Inches(1.0)
            top = Inches(1.5)  # 标题下方开始，可以调整
            width = Inches(8.0)  # 使用较宽的宽度
            height = Inches(4.5)  # 给足够的高度用于垂直居中，可以调整

            text_box_content = shapes.add_textbox(left, top, width, height)
            tf = text_box_content.text_frame

            tf.clear()  # Clear existing content
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 意味着如果文本内容超出文本框的固定大小时，会自动缩小字体大小以适应文本框。
            tf.word_wrap = True  # 启用自动换行

            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # <--- 垂直居中

            paragraph = tf.paragraphs[0]
            paragraph.text = content_text
            self.processing_md_str(self.content.replace("<p>", "").replace("</p>", "\n"))
            self.get_font(paragraph.font, MarkdownCategory.CONTENT)

            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # <--- 水平居中
        else:
            if self.template is None:
                if has_pic:
                    # 下面的代码不引入组件
                    # left = Inches(0.3)
                    # top = Inches(1.5)
                    # width = Inches(5)
                    # height = Inches(5.5)
                    # text_box_content = shapes.add_textbox(left, top, width, height)
                    # tf = text_box_content.text_frame
                    #
                    # tf.clear()  # Clear existing content
                    # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE   # 意味着如果文本内容超出文本框的固定大小时，会自动缩小字体大小以适应文本框。
                    # tf.word_wrap = True  # 启用自动换行
                    #
                    # # 设置文本框架的垂直对齐方式
                    # tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # 普通页面顶部对齐
                    #
                    # # 添加正文
                    # paragraph = tf.paragraphs[0]
                    # paragraph.text = content_text
                    # self.processing_md_str(self.content.replace("<p>", "").replace("</p>", "\n"))
                    # self.get_font(paragraph.font, MarkdownCategory.CONTENT)
                    #
                    # # 设置文本在文本框的水平对齐方式
                    # paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    agent = ComponentAgentHalf()
                    slide_elements_data_positions = agent.get_layout_plan(content_text)
                    # agent = SplitComponentAgent()
                    # slide_elements_data_positions = agent.generate_slide_layout_from_markdown(content_text)
                    agent_type_used = "ComponentAgent"
                    data_to_log = slide_elements_data_positions
                    logging.info(f"Agent used: {agent_type_used}")
                    logging.info(f"slide_elements_data_positions: \n{data_to_log}")
                    self.create_slide_from_elements_content(slide_elements_data_positions, self.infra_color)

                else:
                    agent = ComponentAgent()
                    slide_elements_data_positions = agent.get_layout_plan(content_text)
                    # agent = SplitComponentAgent()
                    # slide_elements_data_positions = agent.generate_slide_layout_from_markdown(content_text)
                    agent_type_used = "ComponentAgent"
                    data_to_log = slide_elements_data_positions
                    logging.info(f"Agent used: {agent_type_used}")
                    logging.info(f"slide_elements_data_positions: \n{data_to_log}")
                    self.create_slide_from_elements_content(slide_elements_data_positions, self.infra_color)

            else:
                if has_pic:
                    random_index = random.randint(0, len(self.template_half) - 1)
                    template_current = self.template_half[random_index]
                    template_current_prompt = template_half_prompt[random_index]
                    slide_elements_data_positions = template_generation(template_current_prompt, content_text,
                                                                        template_current)
                    self.create_slide_from_elements_content(slide_elements_data_positions, self.infra_color)
                else:
                    # 随机取出一个元素
                    random_index = random.randint(0, len(self.template) - 1)
                    template_current = self.template[random_index]
                    template_current_prompt = template_1_prompt[random_index]
                    slide_elements_data_positions = template_generation(template_current_prompt, content_text, template_current)
                    self.create_slide_from_elements_content(slide_elements_data_positions, self.infra_color)

    def _sanitize_filename(self, name, default_extension=".jpg"):
        """清理字符串以使其成为安全的文件名，并添加扩展名。"""
        # 移除或替换文件名中不允许的字符
        sanitized_name = re.sub(r'[<>:"/\\|?*\s]', '_', name)
        # 限制长度
        sanitized_name = sanitized_name[:50]  # 保持合理的短度
        if not sanitized_name:  # 处理空的 alt 文本
            import uuid
            sanitized_name = str(uuid.uuid4())  # 生成一个唯一名称

        # 确保它有扩展名
        if not os.path.splitext(sanitized_name)[1]:
            sanitized_name += default_extension
        return sanitized_name

    def _get_image_extension_from_url(self, url):
        """尝试从 URL 路径猜测图片扩展名，或默认为 .jpg"""
        try:
            path = requests.utils.urlparse(url).path
            ext = os.path.splitext(path)[1]
            if ext.lower() in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                return ext
        except Exception:
            pass
        return ".jpg"  # 默认扩展名

    def _download_image(self, img_url, alt_text, save_dir="picture"):
        """从 URL 下载图片并将其保存到本地。"""
        if 'DeepPPT/core/data' in img_url:
            img_url = img_url.split('DeepPPT/core/data',1)[1]
            img_url = '../DeepPPT/core/data' + img_url
            return img_url

        if not (img_url.startswith("http://") or img_url.startswith("https://")):
            print(f"跳过非 HTTP/S URL: {img_url}")
            return None

        os.makedirs(save_dir, exist_ok=True)  # 创建图片保存目录，如果不存在的话

        extension = self._get_image_extension_from_url(img_url)
        filename = self._sanitize_filename(alt_text if alt_text else "图片",
                                           default_extension=extension)  # 如果alt_text为空，则使用"图片"
        local_image_path = os.path.join(save_dir, filename)

        # 如果文件已存在且非空，则避免重新下载
        if os.path.exists(local_image_path) and os.path.getsize(local_image_path) > 0:
            return local_image_path

        try:
            print(f"正在从以下 URL 下载图片: {img_url} (alt: '{alt_text}')")
            response = requests.get(img_url, stream=True, timeout=10)
            response.raise_for_status()  # HTTP 错误则引发异常
            with open(local_image_path, 'wb') as f:
                for chunk in response.iter_content(8192):
                    f.write(chunk)
            print(f"图片成功下载到: {local_image_path}")
            return local_image_path
        except requests.exceptions.RequestException as e:
            print(f"下载 {img_url} 时出错: {e}")
        except IOError as e:
            print(f"保存图片到 {local_image_path} 时出错: {e}")
        return None

    def _make_page_left(self, image_local_path, left_in=5.5, top_in=1.2, width_in=4, height_in=None):

        try:
            # 定义位置和初始宽度
            left = Inches(left_in)
            top = Inches(top_in)
            width_ppt = Inches(width_in)
            height_ppt = None  # 高度先设为 None

            slide_height = Inches(6.5) # 幻灯片总高度(EMU)
            max_height = (slide_height - top) # 最大允许高度(EMU)

            if height_in is not None:
                height_ppt = Inches(height_in)  # 如果指定了高度，则使用指定高度
            else:
                # 使用 Pillow 获取图片尺寸以保持宽高比
                try:
                    with Image.open(image_local_path) as img:
                        img_width_px, img_height_px = img.size
                        if img_width_px > 0:  # 避免除以零
                            aspect_ratio = img_height_px / img_width_px
                            height_ppt = width_ppt * aspect_ratio  # 根据宽度和宽高比计算高度
                        else:  # 如果图片宽度为零（对于有效图片不应发生）
                            height_ppt = Inches(3)  # 默认高度
                except Exception as e:
                    print(f"无法使用 Pillow 读取图片 {image_local_path} 的尺寸: {e}。将使用默认高度。")
                    # 如果 Pillow 失败或图片损坏，则使用回退高度
                    height_ppt = Inches(3)  # 如果无法确定宽高比，则使用默认高度

            # 检查图片高度是否超过可用空间
            if height_ppt > max_height:
                # 计算需要缩小的比例
                scale_factor = float(max_height / height_ppt)
                # 按比例调整宽度和高度
                new_height = max_height
                new_width = width_ppt * scale_factor
                width_ppt, height_ppt = new_width, new_height
                # print(f"调整图片大小: 原始尺寸 {width_in}×{height_ppt.inches:.1f} 英寸, "
                      # f"调整为 {width_ppt.inches:.1f}×{height_ppt.inches:.1f} 英寸")

            pic = self.slide.shapes.add_picture(image_local_path, left, top, width_ppt, height_ppt)
            print(f"成功添加图片: {image_local_path}")
            return True

        except Exception as e:
            print(f"无法将图片 {image_local_path} 添加到幻灯片: {e}")
            return False

    def _add_images_from_content(self, content):
        image_pattern = re.compile(r'!\[([^\]]*)\]\((.*?)\)')
        images_were_added_to_slide = False

        # 使用列表收集内容部分，以避免在迭代期间修改字符串时出现问题
        content_parts = []
        last_end = 0

        for match_obj in image_pattern.finditer(content):
            alt, img_url = match_obj.groups()  # 提取 alt 文本和图片 URL
            start, end = match_obj.span()

            content_parts.append(content[last_end:start])
            last_end = end  # 更新下一个非图片内容的起始位置

            local_image_path = self._download_image(img_url, alt)

            if local_image_path:
                # 调用 _make_page_left 将下载的图片添加到幻灯片
                # 你可能希望使位置和大小可配置
                images_were_added_to_slide_tag = self._make_page_left(local_image_path)
                if images_were_added_to_slide_tag is True:
                    images_were_added_to_slide = True
            else:
                print(f"下载或添加图片失败: {img_url}")

        # 添加最后一个图片之后剩余的内容
        content_parts.append(content[last_end:])
        processed_content = "".join(content_parts)  # 合并所有部分得到处理后的内容

        return processed_content, images_were_added_to_slide

    def processing_md_str(self, md_str):
        md = markdown.Markdown()

    def create_slide_from_elements_content(self, elements_data, custom_theme_colors=None):
        """
        根据元素数据列表创建一张幻灯片。
        """

        active_theme_colors = custom_theme_colors if custom_theme_colors else DEFAULT_THEME_COLORS
        try:
            for element in elements_data['layout_elements']:
                el_type = element['element_type']
                params = element.get('content_params', {})
                if 'text_size' in element:
                    text_size = element.get('text_size', {})
                    if not isinstance(text_size, dict):
                        text_size = params.get('text_size', {})
                if 'text_size' in params:
                    text_size = params.get('text_size', {})
                    if not isinstance(text_size, dict):
                        text_size = element.get('text_size', {})
                pos = element['position']
                size = element['size']
                style = element.get('style_id')  # Will be None if not applicable

                left_in = Inches(pos['left'])
                top_in = Inches(pos['top'])
                width_in = Inches(size['width'])
                if 'height' in size:
                    height_in = Inches(size['height'])

                if el_type == 'main_text':
                    add_main_text_box(self.slide,
                                      text_content=params.get('text_content', ''),  # Updated key
                                      left=left_in, top=top_in, width=width_in, height=height_in,
                                      theme_colors=active_theme_colors
                                      )

                elif el_type == 'quote_box':
                    add_styled_quote_box(self.slide,
                                         quote_text=params.get('quote_text', ''),  # Updated key
                                         attribution_text=params.get('attribution_text'),  # Updated key
                                         left=left_in, top=top_in, width=width_in, height=height_in,
                                         style_id=style,
                                         theme_colors=active_theme_colors)

                elif el_type == 'data_callout':
                    if height_in is None:  # 如果数据中未提供高度，则使用默认值或进行估算
                        height_in = Inches(1.5)  # KeyDataCallout 的默认高度

                    add_key_data_callout(self.slide,
                                         statistic_value=params.get('statistic_value', ''),  # Updated key
                                         label_text=params.get('label_text', ''),  # Updated key
                                         left=left_in, top=top_in, width=width_in, height=height_in,
                                         style_id=style,
                                         theme_colors=active_theme_colors,
                                         title_size=text_size['title_size'],
                                         content_size=text_size['content_size'],
                                         )

                elif el_type == 'divider':
                    # 分隔线通常不使用 height，其位置由 top 决定，长度由 width 决定
                    add_decorative_divider(self.slide,
                                           left=left_in, top=top_in, width=width_in,
                                           style_id=style,
                                           theme_colors=active_theme_colors)
                elif el_type == 'icon_title_bullets':
                    add_icon_title_bullets_component(self.slide,
                                                     content_data=params.get('content_data', {}),
                                                     icon_emoji=params.get('icon_emoji', '🚀'),
                                                     # Provide default if missing
                                                     left=left_in, top=top_in, width=width_in, height=height_in,
                                                     theme_colors=active_theme_colors,
                                                     title_size=text_size['title_size'],
                                                     content_size=text_size['content_size'],
                                                     emoji_size=text_size['emoji_size'],
                                                     )
                elif el_type == 'triple_card':
                    add_triple_card_component(self.slide,
                                              cards_data=params.get('cards_data', []),  # Ensure it's a list
                                              card_top_emojis=params.get('card_top_emojis'),  # Can be None
                                              card_icon_chars=params.get('card_icon_chars'),  # Can be None
                                              left=left_in, top=top_in, width=width_in, height=height_in,
                                              title_size=text_size['title_size'],
                                              content_size=text_size['content_size'],
                                              emoji_size=text_size['emoji_size'],
                                              theme_colors=active_theme_colors
                                              )
                elif el_type == 'parallel_boxes':
                    add_parallel_boxes_component(self.slide,
                                                 title_text=params.get('title_text', ''),
                                                 content_text=params.get('content_text', ''),
                                                 emoji_left=params.get('emoji_left'),  # Can be None
                                                 emoji_right=params.get('emoji_right'),  # Can be None
                                                 left=left_in, top=top_in, width=width_in, height=height_in,
                                                 content_size=text_size['content_size'],
                                                 emoji_size=text_size['emoji_size'],
                                                 theme_colors=active_theme_colors
                                                 # Spacing, font sizes defaulted in component function
                                                 )
                else:
                    print(f"未知元素类型: {el_type}")
        except:
            print(elements_data)