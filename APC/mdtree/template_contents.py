from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ---- （1）保留你给定的 DEFAULT_THEME_COLORS，不做修改 ----
DEFAULT_THEME_COLORS = {
    "primary_accent": RGBColor(0, 112, 192),       # 主强调色 (例如蓝色)
    "secondary_accent": RGBColor(0, 176, 80),      # 次强调色 (例如绿色)
    "light_background": RGBColor(240, 240, 240),   # 浅色背景 (例如浅灰色)
    "dark_text": RGBColor(30, 30, 30),             # 深色文本
    "light_text": RGBColor(250, 250, 250),         # 浅色文本 (用于深色背景)
    "quote_border": RGBColor(191, 191, 191),       # 引用框边框颜色
    "data_callout_bg": RGBColor(220, 230, 241),    # 数据标注背景色
    "title_text_color": RGBColor(10, 10, 10),      # 标题文本颜色 (专用)
    "content_text_color": RGBColor(50, 50, 50),    # 内容文本颜色 (专用)
}


def create_styled_toc_slide(slide, toc_text_string, theme_colors=None):
    """
    在给定的幻灯片上创建一个样式化的目录页，使用你提供的 DEFAULT_THEME_COLORS。
    修复了“第一个段落始终空白”的问题，并将装饰性元素（点阵 + 汉堡图标）
    的颜色固定为 secondary_accent，避免与左侧深色面板（primary_accent）重合。

    参数：
        slide: pptx.Slide 对象
        toc_text_string: 多行目录项，形如 '• 项目一\n• 项目二\n...'
                         会自动去掉每行开头的 '• '。
        theme_colors: 可选，若传入 dict 则覆盖 DEFAULT_THEME_COLORS，否则使用后者。
    """

    shapes = slide.shapes
    # 如果调用时传入了 theme_colors，且它是 dict，就用它；否则用默认
    theme_colors = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS

    # ---- （2）解析传入的 toc_text_string，去掉“• ”前缀 ----
    toc_lines = []
    if toc_text_string:
        raw_lines = toc_text_string.strip().split('\n')
        for line in raw_lines:
            cleaned_line = line.lstrip('• \t').strip()
            if cleaned_line:
                toc_lines.append(cleaned_line)

    # ---- （3）从 theme_colors 中提取对应颜色 ----
    left_panel_bg_color = theme_colors["primary_accent"]
    main_title_text_color = theme_colors["light_text"]
    sub_title_text_color = theme_colors["light_text"]
    toc_number_color = theme_colors["primary_accent"]     # 目录数字用主强调色
    toc_item_text_color = theme_colors["dark_text"]
    # 这里把装饰性元素颜色固定为次强调色，避免与左侧深色面板(primary_accent)重复
    decorative_elements_color = theme_colors["secondary_accent"]
    slide_bg_color = theme_colors["light_background"]

    # ---- （4）设置幻灯片整体背景为浅灰 ----
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = slide_bg_color

    # ---- （5）创建左侧深蓝色面板 ----
    left_panel_width = Inches(4.2)  # 和你原来代码一致
    left_panel = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0),
        width=left_panel_width,
        height=Inches(7.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = left_panel_bg_color
    left_panel.line.fill.background()  # 无边框

    # ---- （6）在左侧面板中添加大标题 "目录" ----
    # 用一个宽度 = left_panel_width - 1"（留左右各 0.5" 边距）
    title_box_left = Inches(0.5)
    title_box_width = left_panel_width - Inches(1.0)

    title_main_shape = shapes.add_textbox(
        title_box_left, Inches(2.8), title_box_width, Inches(1.2)
    )
    tf_main_title = title_main_shape.text_frame
    tf_main_title.clear()  # 清空自动生成的第一个段落
    tf_main_title.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p_main_title = tf_main_title.paragraphs[0]
    p_main_title.text = "目录"
    p_main_title.alignment = PP_ALIGN.CENTER
    font_main = p_main_title.font
    font_main.name = "微软雅黑 Bold"      # 可以改为 'Microsoft YaHei UI Bold'
    font_main.size = Pt(60)
    font_main.bold = True
    font_main.color.rgb = main_title_text_color

    # ---- （7）在左侧面板中添加副标题 "CONTENTS" ----
    title_sub_shape = shapes.add_textbox(
        title_box_left, Inches(3.8), title_box_width, Inches(0.5)
    )
    tf_sub_title = title_sub_shape.text_frame
    tf_sub_title.clear()
    tf_sub_title.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    p_sub_title = tf_sub_title.paragraphs[0]
    p_sub_title.text = "CONTENTS"
    p_sub_title.alignment = PP_ALIGN.CENTER
    font_sub = p_sub_title.font
    font_sub.name = "微软雅黑 Bold"
    font_sub.size = Pt(16)
    font_sub.bold = False
    font_sub.color.rgb = sub_title_text_color

    # ---- （8）在右侧添加目录项列表 ----
    toc_list_left = left_panel_width + Inches(0.6)  # 距离左侧面板0.6"
    toc_list_top = Inches(1.5)
    toc_list_width = Inches(10.0) - toc_list_left - Inches(0.5)  # 右侧留 0.5" 边距

    # 预估每行文本占 ~0.75" 高度
    toc_list_height = Inches(len(toc_lines) * 0.75 if toc_lines else 1.0)

    toc_textbox = shapes.add_textbox(
        toc_list_left, toc_list_top, toc_list_width, toc_list_height
    )
    tf_toc = toc_textbox.text_frame
    tf_toc.clear()

    # ---- （8.1）如果有目录行，先用第一个段落写第一行 ----
    if toc_lines:
        first_para = tf_toc.paragraphs[0]
        # 第一行编号
        run_num = first_para.add_run()
        run_num.text = f"{1:02d}"
        font_num = run_num.font
        font_num.name = "微软雅黑 Bold"
        font_num.size = Pt(16)
        font_num.bold = True
        font_num.color.rgb = toc_number_color

        # 加三个空格作为编号与文字间隔
        run_spacer = first_para.add_run()
        run_spacer.text = "   "

        # 第一行文字
        run_txt = first_para.add_run()
        run_txt.text = toc_lines[0]
        font_txt = run_txt.font
        font_txt.name = "微软雅黑 Light"
        font_txt.size = Pt(16)
        font_txt.bold = True
        font_txt.color.rgb = toc_item_text_color

        # ---- （8.2）后续条目用 add_paragraph() ----
        for idx, item in enumerate(toc_lines[1:], start=1):
            p = tf_toc.add_paragraph()
            p.space_before = Pt(10)  # 段前加一些空间

            run_n = p.add_run()
            run_n.text = f"{idx+1:02d}"
            font_n = run_n.font
            font_n.name = "微软雅黑 Bold"
            font_n.size = Pt(16)
            font_n.bold = True
            font_n.color.rgb = toc_number_color

            run_sp = p.add_run()
            run_sp.text = "   "

            run_t = p.add_run()
            run_t.text = item
            font_t = run_t.font
            font_t.name = "微软雅黑 Light"
            font_t.size = Pt(16)
            font_t.bold = True
            font_t.color.rgb = toc_item_text_color
    else:
        # 没有任何目录行，留空
        empty_para = tf_toc.paragraphs[0]
        empty_para.text = ""
        empty_para.font.size = Pt(16)

    # ---- （9）添加右上角汉堡菜单图标（三条线），使用次强调色 ----
    bar_h = Inches(0.04)    # 线条高度
    bar_w = Inches(0.4)     # 线条宽度
    menu_icon_l = Inches(10.0 - bar_w - 0.3)  # 距右边 0.3"
    menu_icon_t_start = Inches(0.4)
    bar_sp = Inches(0.07)   # 行间距

    for i in range(3):
        bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            menu_icon_l,
            menu_icon_t_start + i * (bar_h + bar_sp),
            bar_w,
            bar_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = decorative_elements_color
        bar.line.fill.background()