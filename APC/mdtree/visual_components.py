from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE  # 用于预定义形状
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE  # 用于文本对齐、垂直锚定和自动调整大小

# --- 默认主题颜色 (如果未提供，则用作备用) ---
DEFAULT_THEME_COLORS = {
    "primary_accent": RGBColor(0, 112, 192),  # 主强调色 (例如蓝色)
    "secondary_accent": RGBColor(0, 176, 80),  # 次强调色 (例如绿色)
    "light_background": RGBColor(240, 240, 240),  # 浅色背景 (例如浅灰色)
    "dark_text": RGBColor(30, 30, 30),  # 深色文本
    "light_text": RGBColor(250, 250, 250),  # 浅色文本 (用于深色背景)
    "quote_border": RGBColor(191, 191, 191),  # 引用框边框颜色
    "data_callout_bg": RGBColor(220, 230, 241),  # 数据标注背景色
    "title_text_color": RGBColor(10, 10, 10),  # 标题文本颜色 (专用)
    "content_text_color": RGBColor(50, 50, 50),  # 内容文本颜色 (专用)
}


def add_styled_quote_box(slide, quote_text, attribution_text=None,
                         left=Inches(1), top=Inches(1.5), width=Inches(8), height=None,  # 添加了 height 参数
                         style_id="default_quote", theme_colors=None):
    """
    在幻灯片上添加一个带样式的引用框。
    Height 现在是一个可选参数。如果为 None，则会进行估算。
    """
    # 如果未提供 theme_colors，则使用 DEFAULT_THEME_COLORS
    current_theme = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS

    # 如果未提供高度，则进行估算
    if height is None:
        # 根据文本长度粗略估算高度
        height_val = Inches(0.75 + (len(quote_text) // 80) * 0.4 + (0.25 if attribution_text else 0))
    else:
        height_val = height

    # --- 样式定义 ---
    if style_id == "default_quote":  # 默认引用框样式
        bg_color = current_theme.get("light_background", DEFAULT_THEME_COLORS["light_background"])
        border_color = current_theme.get("quote_border", DEFAULT_THEME_COLORS["quote_border"])
        quote_font_color = current_theme.get("dark_text", DEFAULT_THEME_COLORS["dark_text"])
        attr_font_color = RGBColor(80, 80, 80)  # 署名使用稍浅的深色文本
    elif style_id == "left_accent_bar_quote":  # 左侧强调条引用框样式
        bg_color = current_theme.get("light_background", DEFAULT_THEME_COLORS["light_background"])  # 或特定的极浅色
        accent_bar_color = current_theme.get("primary_accent", DEFAULT_THEME_COLORS["primary_accent"])
        quote_font_color = current_theme.get("dark_text", DEFAULT_THEME_COLORS["dark_text"])
        attr_font_color = RGBColor(80, 80, 80)

        # 为左侧强调条单独绘制一个矩形
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(0.15), height_val
        )
        accent_bar.fill.solid()  # 纯色填充
        accent_bar.fill.fore_color.rgb = accent_bar_color  # 设置填充色
        accent_bar.line.fill.background()  # 强调条本身没有边框

        # 调整主文本框的位置和宽度以容纳强调条
        left += Inches(0.25)
        width -= Inches(0.25)
    else:  # 备用样式 (同默认)
        bg_color = current_theme.get("light_background", DEFAULT_THEME_COLORS["light_background"])
        border_color = current_theme.get("quote_border", DEFAULT_THEME_COLORS["quote_border"])
        quote_font_color = current_theme.get("dark_text", DEFAULT_THEME_COLORS["dark_text"])
        attr_font_color = RGBColor(80, 80, 80)

    # --- 引用框的主形状 ---
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height_val  # 使用圆角矩形以获得更柔和的外观
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color  # 设置背景色

    if style_id == "default_quote":  # 根据此逻辑，只有默认样式有边框
        shape.line.color.rgb = border_color  # 设置边框颜色
        shape.line.width = Pt(1.5)  # 设置边框宽度
    else:  # 其他样式无边框 (left_accent_bar_quote 的主框也没有边框)
        shape.line.fill.background()

    # --- 文本框和内容 ---
    tf = shape.text_frame  # 获取文本框对象
    tf.clear()  # 清除任何默认文本
    # 设置文本框内的边距
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.word_wrap = True  # 确保文本在形状内换行
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 文本垂直居中

    # 添加引用段落并设置样式
    p_quote = tf.add_paragraph()
    p_quote.text = f"“{quote_text}”"  # 使用印刷引号
    p_quote.font.name = "微软雅黑 Light"  # 考虑使字体可配置
    p_quote.font.size = Pt(16)
    p_quote.font.italic = True  # 斜体
    p_quote.font.color.rgb = quote_font_color
    p_quote.alignment = PP_ALIGN.CENTER  # 引用文本居中

    # 如果提供了署名，则添加署名段落并设置样式
    if attribution_text:
        p_attr = tf.add_paragraph()
        p_attr.text = f"— {attribution_text}"  # 署名使用破折号
        p_attr.alignment = PP_ALIGN.RIGHT  # 署名右对齐
        p_attr.font.name = "微软雅黑"
        p_attr.font.size = Pt(12)
        p_attr.font.color.rgb = attr_font_color
        p_attr.space_before = Pt(6)  # 在署名前添加一些间距

    # print(f"已添加样式引用框: '{quote_text[:30]}...' 样式为 '{style_id}'")
    return shape


def add_key_data_callout(slide, statistic_value, label_text,
                         left=Inches(1), top=Inches(2), width=Inches(3), height=Inches(1.5),
                         style_id="default_data", theme_colors=None, title_size=20, content_size=10):
    """
    在幻灯片上添加一个关键数据或统计摘要框。
    """
    current_theme = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS

    title_size = Pt(title_size)
    content_size = Pt(content_size)
    # --- 样式定义 ---
    if style_id == "default_data":  # 默认数据标注样式
        bg_color = current_theme.get("data_callout_bg", DEFAULT_THEME_COLORS["data_callout_bg"])
        stat_color = current_theme.get("primary_accent", DEFAULT_THEME_COLORS["primary_accent"])
        label_color = current_theme.get("dark_text", DEFAULT_THEME_COLORS["dark_text"])
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE  # 圆角矩形
    elif style_id == "circle_emphasis_data":  # 圆形强调数据标注样式
        bg_color = current_theme.get("secondary_accent", DEFAULT_THEME_COLORS["secondary_accent"])
        stat_color = current_theme.get("light_text", DEFAULT_THEME_COLORS["light_text"])
        label_color = current_theme.get("light_text", DEFAULT_THEME_COLORS["light_text"])
        shape_type = MSO_SHAPE.OVAL  # 椭圆 (用于圆形)
        min_dim = min(width, height)  # 确保是圆形，取宽高较小值
        width = height = min_dim
    else:  # 备用样式 (同默认)
        bg_color = current_theme.get("data_callout_bg", DEFAULT_THEME_COLORS["data_callout_bg"])
        stat_color = current_theme.get("primary_accent", DEFAULT_THEME_COLORS["primary_accent"])
        label_color = current_theme.get("dark_text", DEFAULT_THEME_COLORS["dark_text"])
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

    # --- 创建形状 ---
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color  # 设置背景色
    shape.line.fill.background()  # 这些标注框默认没有边框

    # --- 文本框和内容 ---
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 文本垂直居中
    tf.word_wrap = False  # 数据标注框通常不换行 (如果文本过长，可能需要调整或启用换行)

    # 添加统计数值段落并设置样式
    p_stat = tf.add_paragraph()
    p_stat.text = statistic_value
    p_stat.font.name = "微软雅黑 Bold"
    p_stat.font.size = title_size
    p_stat.font.color.rgb = stat_color
    p_stat.alignment = PP_ALIGN.CENTER  # 居中对齐

    # 添加标签文本段落并设置样式
    p_label = tf.add_paragraph()
    p_label.text = label_text
    p_label.font.name = "微软雅黑"
    p_label.font.size = content_size
    p_label.font.color.rgb = label_color
    p_label.alignment = PP_ALIGN.CENTER  # 居中对齐
    p_label.space_before = Pt(3)  # 在标签前添加少量间距

    # print(f"已添加关键数据标注: '{statistic_value} - {label_text}' 样式为 '{style_id}'")
    return shape


def add_decorative_divider(slide,
                           left=Inches(0.5), top=Inches(2), width=Inches(9),
                           style_id="default_line", theme_colors=None):
    """
    在幻灯片上添加一个装饰性的分隔线。
    """
    current_theme = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS
    shape_or_shapes = []  # 用于存储创建的形状
    # 对于直线连接符，类型值通常为 1 (MSO_CONNECTOR.STRAIGHT)。
    # from pptx.enum.connector import MSO_CONNECTOR_TYPE (MSO_CONNECTOR_TYPE.STRAIGHT is 1)
    straight_connector_type = 1

    if style_id == "default_line":  # 默认线条样式
        line_color = current_theme.get("primary_accent", DEFAULT_THEME_COLORS["primary_accent"])
        line_weight = Pt(2.0)  # 线条粗细
        shape = slide.shapes.add_connector(
            straight_connector_type,  # 直线连接符
            left, top,  # 起点
            left + width, top  # 终点 (定义了长度和方向)
        )
        shape.line.color.rgb = line_color  # 设置线条颜色
        shape.line.width = line_weight  # 设置线条粗细
        shape_or_shapes = shape  # 单个形状
    elif style_id == "dots_accent_line":  # 带点的强调线条样式
        # 这种样式会创建一条主线和中间的一个点。
        line_color = current_theme.get("quote_border", DEFAULT_THEME_COLORS["quote_border"])  # 使用更柔和的线条颜色
        dot_color = current_theme.get("secondary_accent", DEFAULT_THEME_COLORS["secondary_accent"])  # 点的颜色
        line_weight = Pt(1)  # 线条粗细

        # 主线
        main_line = slide.shapes.add_connector(
            straight_connector_type,
            left, top,
            left + width, top
        )
        main_line.line.color.rgb = line_color
        main_line.line.width = line_weight
        shape_or_shapes.append(main_line)  # 添加到形状列表

        # 中间的点
        dot_size = Inches(0.1)  # 点的大小
        # 计算点的位置，使其在线条上居中。
        dot_left = left + (width / 2) - (dot_size / 2)
        dot_top = top - (dot_size / 2)  # 使点在垂直方向上与线条对齐
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, dot_left, dot_top, dot_size, dot_size  # 圆形
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color  # 点的填充色
        dot.line.fill.background()  # 点本身没有边框
        shape_or_shapes.append(dot)  # 添加到形状列表
    else:  # 备用样式 (同默认)
        line_color = current_theme.get("primary_accent", DEFAULT_THEME_COLORS["primary_accent"])
        line_weight = Pt(2.0)
        shape = slide.shapes.add_connector(
            straight_connector_type,
            left, top,
            left + width, top
        )
        shape.line.color.rgb = line_color
        shape.line.width = line_weight
        shape_or_shapes = shape

    # print(f"已添加装饰性分隔线，样式为 '{style_id}'")
    return shape_or_shapes  # 返回形状或形状列表


def add_main_text_box(slide, text_content,
                      left, top, width, height,
                      theme_colors=None, font_size_pt=14, font_name="微软雅黑", text_alignment=PP_ALIGN.LEFT):
    """
    在幻灯片上添加一个主文本框。
    """
    current_theme = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS
    text_color = current_theme.get("content_text_color", DEFAULT_THEME_COLORS["content_text_color"])

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()  # 清除默认段落

    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # 主文本通常顶部对齐
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 如果文本过多，缩小字体以适应形状

    # 直接设置文本内容，\n 会被正确处理为换行
    tf.text = text_content

    # 为文本框中的所有段落设置统一样式
    for p in tf.paragraphs:
        p.font.name = font_name
        p.font.size = Pt(font_size_pt)
        p.font.color.rgb = text_color
        p.alignment = text_alignment
        # 可以根据需要设置行距等
        # p.line_spacing = 1.0 # 例如，单倍行距

    # print(f"已添加主文本框: '{text_content[:30]}...'")
    return textbox


def add_icon_title_bullets_component(slide, content_data,
                                     left, top, width, height,  # MODIFIED: 组件的整体边界框
                                     icon_emoji="🚀",
                                     title_size=20,
                                     content_size=14,
                                     emoji_size=72,
                                     theme_colors=None):  # Will be used now

    """
    添加一个组件，该组件左侧带有一个圆角图标框（内含Emoji和标题），右侧带有项目符号点文本区。

    参数：
        slide (python-pptx幻灯片对象)：python-pptx 的幻灯片对象。
        content_data (dict)：一个字典，包含 'title' (str) 和 'bullets' (list of str)。
                             例如: {'title': "组件标题", 'bullets': ["要点1", "要点2"]}
        left (Emu): 组件左上角的 X 坐标。如果为 None，则基于 slide_width 水平居中。
        top (Emu): 组件左上角的 Y 坐标。如果为 None，则基于 slide_height 垂直居中。
        width (Emu): 组件的总宽度。
        height (Emu): 组件的总高度 (这也将是图标框和项目符号区域的高度)。
        icon_emoji (str)：左侧图标框中使用的表情符号字符。
        font_name (str, 可选): 文本使用的字体名称。默认 "微软雅黑"。
        emoji_size (int, 可选)：表情符号的字体大小（单位：磅）。默认 72。
        title_size (int, 可选)：图标框内标题文本的字体大小（单位：磅）。默认 20。
        content_size (int, 可选)：项目符号点文本的字体大小（单位：磅）。默认 14。
        theme_colors (dict, 可选): 自定义主题颜色字典。如果未提供，则使用 DEFAULT_THEME_COLORS。
    """
    # --- Theme Color Setup ---
    current_theme = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS

    def get_color_from_theme(key_name, default_colors_dict):
        base_color = default_colors_dict.get(key_name)
        return current_theme.get(key_name, base_color)

    icon_box_bg_color = get_color_from_theme("primary_accent", DEFAULT_THEME_COLORS)
    icon_box_text_color = get_color_from_theme("light_text", DEFAULT_THEME_COLORS)
    bullets_area_bg_color = get_color_from_theme("light_background", DEFAULT_THEME_COLORS)
    bullet_char_custom_color = get_color_from_theme("primary_accent", DEFAULT_THEME_COLORS)  # 项目符号颜色
    bullet_text_custom_color = get_color_from_theme("dark_text", DEFAULT_THEME_COLORS)  # 项目文本颜色

    font_name = "微软雅黑"  # RENAMED from test_front_name
    icon_box_width_ratio = 0.22  # ADDED: 图标框宽度占总宽度的比例
    spacing_between_elements = Inches(0.1)  # ADDED: 图标框和文本区之间的间距
    box_roundness = 0.1
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    # --- Extract data from content_data ---
    title_text = content_data.get('title', "Default Title")
    bullet_points = content_data.get('bullets', ["Default bullet point."])

    # --- Component Dimensions from parameters ---
    component_total_width = width
    component_total_height = height

    calculated_icon_box_width = component_total_width * icon_box_width_ratio
    calculated_bullets_area_width = component_total_width - calculated_icon_box_width - spacing_between_elements

    if calculated_icon_box_width <= 0 or calculated_bullets_area_width <= 0:
        raise ValueError("Calculated width for icon box or bullets area is too small or negative. "
                         "Adjust component width, icon_box_width_ratio, or spacing_between_elements.")

    # --- Positioning ---
    actual_left = left
    actual_top = top

    if actual_left is None:
        if slide_width is None:
            raise ValueError("If 'left' is None, 'slide_width' must be provided for horizontal centering.")
        actual_left = (slide_width - component_total_width) / 2
    if actual_top is None:
        if slide_height is None:
            raise ValueError("If 'top' is None, 'slide_height' must be provided for vertical centering.")
        actual_top = (slide_height - component_total_height) / 2

    created_shapes = []

    # --- 1. Left Icon Box ---
    icon_box_l_pos = actual_left
    icon_box_t_pos = actual_top

    shape_icon_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        icon_box_l_pos, icon_box_t_pos,
        calculated_icon_box_width, component_total_height  # Use component_total_height for box height
    )
    shape_icon_box.fill.solid()
    shape_icon_box.fill.fore_color.rgb = icon_box_bg_color
    shape_icon_box.line.fill.background()  # No border
    try:
        shape_icon_box.adjustments[0] = box_roundness
    except IndexError:
        pass
    created_shapes.append(shape_icon_box)

    tf_icon_box = shape_icon_box.text_frame
    tf_icon_box.clear()
    tf_icon_box.margin_left = Inches(0.1)
    tf_icon_box.margin_right = Inches(0.1)
    tf_icon_box.margin_top = Inches(0.15)  # Adjusted for potentially tighter vertical space
    tf_icon_box.margin_bottom = Inches(0.1)
    tf_icon_box.word_wrap = True
    tf_icon_box.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Try Middle anchor for better emoji/title balance

    if icon_emoji and icon_emoji.strip():
        p_emoji = tf_icon_box.add_paragraph()
        run_emoji = p_emoji.add_run()
        run_emoji.text = icon_emoji
        run_emoji.font.name = font_name
        run_emoji.font.size = Pt(emoji_size)
        run_emoji.font.color.rgb = icon_box_text_color
        p_emoji.alignment = PP_ALIGN.CENTER
        # Adjust space_after based on whether title is also present
        if title_text and title_text.strip():
            p_emoji.space_after = Pt(3)  # Smaller space if title follows
        else:
            p_emoji.space_after = Pt(0)

    if title_text and title_text.strip():
        p_title = tf_icon_box.add_paragraph()
        run_title = p_title.add_run()
        run_title.text = title_text
        run_title.font.name = font_name
        run_title.font.size = Pt(title_size)
        run_title.font.color.rgb = icon_box_text_color
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_before = Pt(0)  # No extra space before title if emoji is right above

    # --- 2. Right Bullets Area ---
    bullets_area_l_pos = icon_box_l_pos + calculated_icon_box_width + spacing_between_elements
    bullets_area_t_pos = actual_top

    shape_bullets_area = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        bullets_area_l_pos, bullets_area_t_pos,
        calculated_bullets_area_width, component_total_height  # Use component_total_height
    )
    shape_bullets_area.fill.solid()
    shape_bullets_area.fill.fore_color.rgb = bullets_area_bg_color
    shape_bullets_area.line.fill.background()  # No border
    try:
        shape_bullets_area.adjustments[0] = box_roundness
    except IndexError:
        pass
    created_shapes.append(shape_bullets_area)

    tf_bullets = shape_bullets_area.text_frame
    tf_bullets.clear()
    tf_bullets.margin_left = Inches(0.25)
    tf_bullets.margin_right = Inches(0.15)
    tf_bullets.margin_top = Inches(0.2)  # Adjusted for better vertical centering of bullets
    tf_bullets.margin_bottom = Inches(0.15)
    tf_bullets.word_wrap = True
    tf_bullets.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # Bullets usually start from top

    for point_text in bullet_points:
        p_bullet = tf_bullets.add_paragraph()
        p_bullet.level = 0  # Ensure it's a top-level bullet

        # Bullet character run
        bullet_char_run = p_bullet.add_run()
        bullet_char_run.text = "• "  # Standard bullet character with a space
        bullet_char_run.font.name = font_name
        bullet_char_run.font.size = Pt(content_size + 2)  # Slightly larger bullet
        bullet_char_run.font.color.rgb = bullet_char_custom_color

        # Bullet text run
        text_run = p_bullet.add_run()
        text_run.text = point_text
        text_run.font.name = font_name
        text_run.font.size = Pt(content_size)
        text_run.font.color.rgb = bullet_text_custom_color

        p_bullet.alignment = PP_ALIGN.LEFT
        p_bullet.space_after = Pt(6)  # Space between bullet points


def add_triple_card_component(slide, cards_data,
                              left, top, width, height,
                              card_top_emojis=None,  # 卡片顶部的 Emoji 列表
                              card_icon_chars=None,  # 新增: 卡片右下角的图标字符列表
                              title_size=16,
                              content_size=10,
                              emoji_size=28,  # 用于 card_top_emojis
                              theme_colors=None):
    """
    在幻灯片上添加一个由三个并行卡片组成的组件。
    卡片顶部可以包含一个可选的居中emoji (通过 card_top_emojis 参数传入)。
    卡片右下角可以包含一个可选的图标字符 (通过 card_icon_chars 参数传入)。

    Args:
        slide: The python-pptx slide object.
        cards_data (list): 包含三个字典的列表，每个字典包含 'title' 和 'content'.
                           例如: [{'title': "T1", 'content': "C1"}, ...]
                           注意：此函数强制要求 cards_data 包含三个元素。
        left (Emu): 组件左上角的X坐标。如果为 None，则组件将基于 slide_width 水平居中。
        top (Emu): 组件左上角的Y坐标。如果为 None，则组件将基于 slide_height 垂直居中。
        width (Emu): 组件的总宽度。
        height (Emu): 组件的总高度 (这也将是每个卡片的高度)。
        card_top_emojis (list, optional): 一个包含最多三个字符串元素的列表，代表卡片顶部的emoji。
                                          例如: ['🎨', '🚀', '💡']。默认为 None。
        card_icon_chars (list, optional): 一个包含最多三个字符串元素的列表，代表卡片右下角的图标字符。
                                          例如: ['📊', '🎯', '❤️']。默认为 None。
        font_name (str, optional): 字体名称。
        emoji_size (int, optional): 卡片顶部主emoji的字体大小 (Pt)。
        title_size (int, optional): 标题字体大小 (Pt)。
        content_size (int, optional): 内容字体大小 (Pt)。
        icon_font_size_pt (int, optional): 右下角图标字符的字体大小 (Pt)。
        theme_colors (dict, optional): 用户可以提供一个主题颜色字典。

    """
    if len(cards_data) != 3:
        raise ValueError("对于 'add_triple_card_component' 函数, 'cards_data' 列表必须包含三个元素。")
    num_cards = 3
    font_name = "微软雅黑"
    if card_top_emojis is None:
        card_top_emojis = ["💡", "🚀", "🌟"]
    if card_icon_chars is None:
        card_icon_chars = ["📈", "📢", "👍"]
    spacing = Inches(0.25)
    current_theme = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    icon_font_size_pt = emoji_size - 4
    # --- 颜色定义 ---
    color_lr_bg_key = "light_text"
    color_lr_text_key = "dark_text"
    color_lr_icon_key = "primary_accent"  # 用于右下角图标颜色 (如果卡片是左右卡片)

    color_mid_bg_dark_key = "primary_accent"
    color_mid_bg_light_strip_key = "data_callout_bg"
    color_mid_text_key = "light_text"
    color_mid_icon_key = "light_text"  # 用于右下角图标颜色 (如果卡片是中间卡片)

    def get_color_from_theme(key_name, default_colors_dict):
        base_color = default_colors_dict.get(key_name)
        return current_theme.get(key_name, base_color)

    color_left_right_bg = get_color_from_theme(color_lr_bg_key, DEFAULT_THEME_COLORS)
    color_left_right_text = get_color_from_theme(color_lr_text_key, DEFAULT_THEME_COLORS)
    color_left_right_icon_color = get_color_from_theme(color_lr_icon_key, DEFAULT_THEME_COLORS)  # 修改变量名

    color_middle_bg_dark = get_color_from_theme(color_mid_bg_dark_key, DEFAULT_THEME_COLORS)
    color_middle_bg_light_strip = get_color_from_theme(color_mid_bg_light_strip_key, DEFAULT_THEME_COLORS)
    color_middle_text = get_color_from_theme(color_mid_text_key, DEFAULT_THEME_COLORS)
    color_middle_icon_color = get_color_from_theme(color_mid_icon_key, DEFAULT_THEME_COLORS)  # 修改变量名

    # --- 组件和卡片尺寸计算 ---
    component_total_width = width
    component_card_height = height

    actual_left = left
    actual_top = top

    if actual_left is None:
        if slide_width is None:
            raise ValueError("如果 'left' 为 None，则必须提供 'slide_width' 以进行水平居中。")
        actual_left = (slide_width - component_total_width) / 2

    if actual_top is None:
        if slide_height is None:
            raise ValueError("如果 'top' 为 None，则必须提供 'slide_height' 以进行垂直居中。")
        actual_top = (slide_height - component_card_height) / 2

    calculated_card_width = (component_total_width - (num_cards - 1) * spacing) / num_cards

    if calculated_card_width <= 0:
        raise ValueError("计算出的卡片宽度过小或为负。请调整组件宽度 'width' 或间距 'spacing'。")

    created_shapes = []

    for i in range(num_cards):
        card_data = cards_data[i]  # 现在只包含 title 和 content
        current_card_left = actual_left + i * (calculated_card_width + spacing)
        current_card_top = actual_top
        is_middle_card = (i == 1)

        card_bg_color_eff = color_middle_bg_dark if is_middle_card else color_left_right_bg
        text_color_eff = color_middle_text if is_middle_card else color_left_right_text
        # 有效的右下角图标颜色取决于卡片位置
        icon_char_color_eff = color_middle_icon_color if is_middle_card else color_left_right_icon_color

        # --- 创建卡片形状 ---
        shape_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            current_card_left, current_card_top, calculated_card_width, component_card_height
        )
        shape_card.fill.solid()
        shape_card.fill.fore_color.rgb = card_bg_color_eff
        shape_card.line.fill.background()
        try:
            shape_card.adjustments[0] = 0.1
        except IndexError:
            pass
        created_shapes.append(shape_card)

        # --- 中间卡片的特殊顶部条带 ---
        strip_height_val = Inches(0)
        if is_middle_card:
            strip_height_val = Inches(0.75)
            strip_height_val = min(strip_height_val, component_card_height)
            shape_strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                current_card_left, current_card_top, calculated_card_width, strip_height_val
            )
            shape_strip.fill.solid()
            shape_strip.fill.fore_color.rgb = color_middle_bg_light_strip
            shape_strip.line.fill.background()
            created_shapes.append(shape_strip)

        # --- 卡片内文本框设置 ---
        tf = shape_card.text_frame
        tf.clear()
        text_margin_x = Inches(0.25)
        text_margin_y_top_base = Inches(0.15)
        tf.margin_left = text_margin_x
        tf.margin_right = text_margin_x
        tf.margin_top = strip_height_val + text_margin_y_top_base if is_middle_card else text_margin_y_top_base
        tf.margin_bottom = Inches(0.5)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        # --- 卡片顶部 Emoji (从 card_top_emojis 参数获取) ---
        card_emoji_char = None
        if card_top_emojis and i < len(card_top_emojis):
            card_emoji_char = card_top_emojis[i]

        if card_emoji_char and card_emoji_char.strip():
            # p_card_emoji = tf.add_paragraph()
            # p_card_emoji.text = card_emoji_char
            # p_card_emoji.font.name = font_name
            # p_card_emoji.font.size = Pt(emoji_size)
            # p_card_emoji.font.color.rgb = text_color_eff # 顶部 emoji 颜色通常跟随卡片主文本颜色
            # p_card_emoji.alignment = PP_ALIGN.CENTER
            # p_card_emoji.space_before = Inches(0.05)
            # p_card_emoji.space_after = Pt(3)
            if is_middle_card:
                # 对于中间卡片，如果条带具有正高度，则 Emoji 放置在条带中
                # 定义 Emoji 文本框的边界，使其与条带相同
                emoji_textbox_left = current_card_left
                emoji_textbox_top = current_card_top  # 条带位于卡片的顶部
                emoji_textbox_width = calculated_card_width
                emoji_textbox_height = strip_height_val

                # 专门为条带上的 Emoji 创建一个文本框
                shape_emoji_strip_textbox = slide.shapes.add_textbox(
                    emoji_textbox_left, emoji_textbox_top,
                    emoji_textbox_width, emoji_textbox_height
                )
                tf_emoji_strip = shape_emoji_strip_textbox.text_frame
                tf_emoji_strip.clear()
                tf_emoji_strip.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 在条带中垂直居中 Emoji
                tf_emoji_strip.margin_left = Inches(0)  # 无内部边距以实现精确对齐
                tf_emoji_strip.margin_right = Inches(0)
                tf_emoji_strip.margin_top = Inches(0)
                tf_emoji_strip.margin_bottom = Inches(0)
                tf_emoji_strip.word_wrap = False  # Emoji 不应换行

                p_emoji_strip = tf_emoji_strip.add_paragraph()
                p_emoji_strip.text = card_emoji_char
                p_emoji_strip.font.name = font_name
                p_emoji_strip.font.size = Pt(emoji_size)
                # 对于中间卡片，text_color_eff 是 color_middle_text。
                # 应选择此颜色以与 color_middle_bg_light_strip 形成良好对比。
                p_emoji_strip.font.color.rgb = text_color_eff
                p_emoji_strip.alignment = PP_ALIGN.CENTER  # 在条带中水平居中 Emoji
                created_shapes.append(shape_emoji_strip_textbox)
            # 如果中间卡片的 strip_height_val 为 0，则 Emoji 不会添加到条带中。
            # 在这种情况下，中间卡片的主文本框 (tf) 中不会添加 Emoji 段落。
            else:
                # 对于非中间卡片，Emoji 像之前一样进入主文本框 (tf)
                p_card_emoji = tf.add_paragraph()
                p_card_emoji.text = card_emoji_char
                p_card_emoji.font.name = font_name
                p_card_emoji.font.size = Pt(emoji_size)
                # 对于这些卡片，text_color_eff 将是 color_left_right_text
                p_card_emoji.font.color.rgb = text_color_eff
                p_card_emoji.alignment = PP_ALIGN.CENTER
                p_card_emoji.space_before = Inches(0.05)
                p_card_emoji.space_after = Pt(3)

        # --- 标题 ---
        p_title = tf.add_paragraph()
        p_title.text = card_data.get('title', "Default Title")  # 从 card_data 获取
        p_title.font.name = font_name
        p_title.font.size = Pt(title_size)
        p_title.font.bold = True
        p_title.font.color.rgb = text_color_eff
        p_title.alignment = PP_ALIGN.LEFT
        p_title.space_after = Pt(2)

        # --- 内容 ---
        p_content = tf.add_paragraph()
        p_content.text = card_data.get('content', "Default content text.")  # 从 card_data 获取
        p_content.font.name = font_name
        p_content.font.size = Pt(content_size)
        p_content.font.color.rgb = text_color_eff
        p_content.alignment = PP_ALIGN.LEFT
        p_content.line_spacing = 1.2

        # --- 右下角图标 (从 card_icon_chars 参数获取) ---
        icon_char_val = None
        if card_icon_chars and i < len(card_icon_chars):
            icon_char_val = card_icon_chars[i]

        if icon_char_val and icon_char_val.strip():
            icon_box_size = Inches(0.5)
            icon_offset_from_edge = Inches(0.15)
            icon_left_pos = current_card_left + calculated_card_width - icon_box_size - icon_offset_from_edge
            icon_top_pos = current_card_top + component_card_height - icon_box_size - icon_offset_from_edge

            shape_icon = slide.shapes.add_textbox(icon_left_pos, icon_top_pos, icon_box_size, icon_box_size)
            tf_icon = shape_icon.text_frame
            tf_icon.clear()
            tf_icon.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            tf_icon.margin_left = tf_icon.margin_right = tf_icon.margin_top = tf_icon.margin_bottom = Inches(0)

            p_icon = tf_icon.add_paragraph()
            p_icon.text = icon_char_val
            p_icon.font.name = font_name
            p_icon.font.size = Pt(icon_font_size_pt)
            p_icon.font.bold = True
            p_icon.font.color.rgb = icon_char_color_eff  # 使用上面确定的有效图标颜色
            p_icon.alignment = PP_ALIGN.CENTER
            created_shapes.append(shape_icon)


def add_parallel_boxes_component(slide, title_text, content_text,
                                 left, top, width, height,  # MODIFIED: Component's overall bounding box
                                 emoji_left="✒️", emoji_right="📄",
                                 content_size=12,
                                 emoji_size=24,
                                 theme_colors=None):
    """
    在幻灯片上添加一个由两个并行框组成的组件。
    Emoji(如果提供)将显示在框的顶部，后跟文本。
    组件的位置、Emoji图案及字体大小可以通过参数指定。

    Args:
        slide: The python-pptx slide object to add the component to.
        title_text: The text for the left box.
        content_text: The text for the right box.
        left (Emu): 组件左上角的X坐标。如果为None，则基于slide_width水平居中。
        top (Emu): 组件左上角的Y坐标。如果为None，则基于slide_height垂直居中。
        width (Emu): 组件的总宽度。
        height (Emu): 组件的总高度 (这也将是每个框的高度)。
        emoji_left (str, optional): 左边框的emoji字符。默认为 "✒️"。传None或空字符串则不显示。
        emoji_right (str, optional): 右边框的emoji字符。默认为 "📄"。传None或空字符串则不显示。
        spacing (Emu, optional): 两个框之间的间距。默认为 Inches(0.2)。
        font_name (str, optional): 主要文本和Emoji的字体名称。默认为 "微软雅黑"。
        emoji_size (int, optional): Emoji的字体大小 (单位 Pt)。默认 24。
        content_size (int, optional): 主要文本的字体大小 (单位 Pt)。默认 12。
        theme_colors (dict, optional): 主题颜色字典。如果未提供，则使用 DEFAULT_THEME_COLORS。
    """
    # --- Theme Color Setup ---
    current_theme = theme_colors if isinstance(theme_colors, dict) else DEFAULT_THEME_COLORS

    font_name = "微软雅黑"

    def get_color_from_theme(key_name, default_colors_dict):
        base_color = default_colors_dict.get(key_name)
        return current_theme.get(key_name, base_color)

    spacing = Inches(0.2)
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    left_box_bg_color = get_color_from_theme("primary_accent", DEFAULT_THEME_COLORS)
    left_box_text_color = get_color_from_theme("light_text", DEFAULT_THEME_COLORS)
    right_box_bg_color = get_color_from_theme("light_background", DEFAULT_THEME_COLORS)
    right_box_text_color = get_color_from_theme("dark_text", DEFAULT_THEME_COLORS)

    # --- Component and Box Dimensions ---
    component_total_width = width
    component_box_height = height  # Each box will take the full height of the component

    # Calculate individual box width
    calculated_box_width = (component_total_width - spacing) / 2
    if calculated_box_width <= 0:
        raise ValueError("Calculated box width is too small or negative. Adjust component width or spacing.")

    # --- Positioning ---
    actual_left = left
    actual_top = top

    if actual_left is None:
        if slide_width is None:
            raise ValueError("If 'left' is None, 'slide_width' must be provided for horizontal centering.")
        slide_center_x = slide_width / 2
        actual_left = slide_center_x - (component_total_width / 2)

    if actual_top is None:
        if slide_height is None:
            raise ValueError("If 'top' is None, 'slide_height' must be provided for vertical centering.")
        slide_center_y = slide_height / 2
        actual_top = slide_center_y - (component_box_height / 2)

    created_shapes = []

    # --- 1. Left Box ---
    left_box_l_pos = actual_left
    left_box_t_pos = actual_top

    shape_left = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left_box_l_pos, left_box_t_pos,
        calculated_box_width, component_box_height
    )
    shape_left.fill.solid()
    shape_left.fill.fore_color.rgb = left_box_bg_color
    shape_left.line.fill.background()  # No border
    try:
        shape_left.adjustments[0] = 0.15  # Rounded corner adjustment
    except IndexError:
        pass  # Shape might not have adjustments
    created_shapes.append(shape_left)

    tf_left = shape_left.text_frame
    tf_left.clear()
    tf_left.margin_left = Inches(0.2)  # Increased margin slightly
    tf_left.margin_right = Inches(0.2)  # Increased margin slightly
    tf_left.margin_top = Inches(0.15)
    tf_left.margin_bottom = Inches(0.15)
    tf_left.word_wrap = True
    tf_left.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    if emoji_left and emoji_left.strip():
        p_emoji_left = tf_left.add_paragraph()
        p_emoji_left.text = emoji_left
        p_emoji_left.font.name = font_name
        p_emoji_left.font.size = Pt(emoji_size)
        p_emoji_left.font.color.rgb = left_box_text_color
        p_emoji_left.alignment = PP_ALIGN.CENTER
        p_emoji_left.space_after = Pt(5)  # Space after emoji

    p_left = tf_left.add_paragraph()
    p_left.text = title_text
    p_left.font.name = font_name
    p_left.font.size = Pt(content_size)
    p_left.font.color.rgb = left_box_text_color
    p_left.alignment = PP_ALIGN.LEFT
    p_left.line_spacing = 1.15  # Slightly more line spacing for readability

    # --- 2. Right Box ---
    right_box_l_pos = left_box_l_pos + calculated_box_width + spacing
    right_box_t_pos = actual_top

    shape_right = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_box_l_pos, right_box_t_pos,
        calculated_box_width, component_box_height
    )
    shape_right.fill.solid()
    shape_right.fill.fore_color.rgb = right_box_bg_color
    shape_right.line.fill.background()  # No border
    try:
        shape_right.adjustments[0] = 0.15  # Rounded corner adjustment
    except IndexError:
        pass
    created_shapes.append(shape_right)

    tf_right = shape_right.text_frame
    tf_right.clear()
    tf_right.margin_left = Inches(0.2)  # Increased margin slightly
    tf_right.margin_right = Inches(0.2)  # Increased margin slightly
    tf_right.margin_top = Inches(0.15)
    tf_right.margin_bottom = Inches(0.15)
    tf_right.word_wrap = True
    tf_right.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    if emoji_right and emoji_right.strip():
        p_emoji_right = tf_right.add_paragraph()
        p_emoji_right.text = emoji_right
        p_emoji_right.font.name = font_name
        p_emoji_right.font.size = Pt(emoji_size)
        p_emoji_right.font.color.rgb = right_box_text_color  # Use right box text color for its emoji
        p_emoji_right.alignment = PP_ALIGN.CENTER
        p_emoji_right.space_after = Pt(5)  # Space after emoji

    p_right = tf_right.add_paragraph()
    p_right.text = content_text
    p_right.font.name = font_name
    p_right.font.size = Pt(content_size)
    p_right.font.color.rgb = right_box_text_color
    p_right.alignment = PP_ALIGN.LEFT
    p_right.line_spacing = 1.15  # Slightly more line spacing for readability


def create_slide_from_elements(prs, slide_layout, elements_data, custom_theme_colors=None):
    """
    根据元素数据列表创建一张幻灯片。
    """
    slide = prs.slides.add_slide(slide_layout)
    active_theme_colors = custom_theme_colors if custom_theme_colors else DEFAULT_THEME_COLORS

    for element in elements_data:
        el_type = element['element_type']
        params = element.get('content_params', {})
        if 'text_size' in element:
            text_size = element.get('text_size', {})
            if not isinstance(text_size, dict):
                text_size = params.get('text_size', {})
        if 'text_size' in params:
            text_size = params.get('text_size', {})
            if not isinstance(text_size, dict):
                text_size = element.get('text_size', {})

        pos = element['position']
        size = element['size']
        style = element.get('style_id')  # Will be None if not applicable

        # 将字典中的单位（假设为英寸）转换为 Inches 对象
        left_in = Inches(pos['left'])
        top_in = Inches(pos['top'])
        width_in = Inches(size['width'])
        if 'height' in size:
            height_in = Inches(size['height'])

        if el_type == 'main_text':
            add_main_text_box(slide,
                              text_content=params.get('text_content', ''),  # Updated key
                              left=left_in, top=top_in, width=width_in, height=height_in,
                              theme_colors=active_theme_colors
                              )
        elif el_type == 'quote_box':
            add_styled_quote_box(slide,
                                 quote_text=params.get('quote_text', ''),  # Updated key
                                 attribution_text=params.get('attribution_text'),  # Updated key
                                 left=left_in, top=top_in, width=width_in, height=height_in,
                                 style_id=style,
                                 theme_colors=active_theme_colors)
        elif el_type == 'data_callout':
            add_key_data_callout(slide,
                                 statistic_value=params.get('statistic_value', ''),  # Updated key
                                 label_text=params.get('label_text', ''),  # Updated key
                                 left=left_in, top=top_in, width=width_in, height=height_in,
                                 style_id=style,
                                 theme_colors=active_theme_colors,
                                 title_size=text_size['title_size'],
                                 content_size=text_size['content_size'],
                                 )
        elif el_type == 'divider':
            # 分隔线使用 top_in 来定位，其视觉高度由内部样式决定
            add_decorative_divider(slide,
                                   left=left_in, top=top_in, width=width_in,
                                   style_id=style,
                                   theme_colors=active_theme_colors)
        elif el_type == 'icon_title_bullets':
            add_icon_title_bullets_component(slide,
                                             content_data=params.get('content_data', {}),
                                             icon_emoji=params.get('icon_emoji', '🚀'),  # Provide default if missing
                                             left=left_in, top=top_in, width=width_in, height=height_in,
                                             theme_colors=active_theme_colors,
                                             title_size=text_size['title_size'],
                                             content_size=text_size['content_size'],
                                             emoji_size=text_size['emoji_size'],
                                             # Optional params like icon_box_width_ratio, font sizes etc.,
                                             # can be added here if LLM is to control them.
                                             # For now, they use defaults in the component function.
                                             )
        elif el_type == 'triple_card':
            add_triple_card_component(slide,
                                      cards_data=params.get('cards_data', []),  # Ensure it's a list
                                      card_top_emojis=params.get('card_top_emojis'),  # Can be None
                                      card_icon_chars=params.get('card_icon_chars'),  # Can be None
                                      left=left_in, top=top_in, width=width_in, height=height_in,
                                      title_size=text_size['title_size'],
                                      content_size=text_size['content_size'],
                                      emoji_size=text_size['emoji_size'],
                                      theme_colors=active_theme_colors
                                      # Font sizes are defaulted in component function
                                      )
        elif el_type == 'parallel_boxes':
            add_parallel_boxes_component(slide,
                                         title_text=params.get('title_text', ''),
                                         content_text=params.get('content_text', ''),
                                         emoji_left=params.get('emoji_left'),  # Can be None
                                         emoji_right=params.get('emoji_right'),  # Can be None
                                         left=left_in, top=top_in, width=width_in, height=height_in,
                                         content_size=text_size['content_size'],
                                         emoji_size=text_size['emoji_size'],
                                         theme_colors=active_theme_colors
                                         # Spacing, font sizes defaulted in component function
                                         )
        else:
            print(f"未知元素类型: {el_type}")

    return slide


# --- 示例数据和调用 ---
if __name__ == '__main__':
    example_elements = \
        [{'element_type': 'quote_box',
          'content_params': {'quote_text': '',
                             'attribution_text': ''}, 'style_id': 'left_accent_bar_quote',
          'position': {'left': 0.5, 'top': 1.5}, 'size': {'width': 9.0, 'height': 1.5}},
         {'element_type': 'triple_card', 'content_params': {
             'cards_data': [{'title': '', 'content': ""},
                            {'title': '', 'content': ""},
                            {'title': '', 'content': ""}],
             'card_top_emojis': ['🚀', '☁️', '🕊️']}, 'position': {'left': 0.5, 'top': 3.25},
          'size': {'width': 9.0, 'height': 3.0},
          'text_size': {'title_size': 20, 'content_size': 12, 'emoji_size': 24}},
         {'element_type': 'main_text',
          'content_params': {'text_content': ''},
          'position': {'left': 0.5, 'top': 6.5}, 'size': {'width': 9.0, 'height': 0.5}}]
    # 创建一个新的演示文稿。
    prs = Presentation()
    # 选择一个空白幻灯片版式 (版式索引 5 或 6 通常是空白，根据模板可能不同)。
    # 索引 5: Blank
    # 索引 6: Title and Content (or sometimes also Blank)
    # It's safer to iterate and find a truly blank one if possible, or just use a common index.
    blank_slide_layout = prs.slide_layouts[5]

    # 使用修正后的数据和布局创建幻灯片
    create_slide_from_elements(prs, blank_slide_layout, example_elements, DEFAULT_THEME_COLORS)


    # (可选) 添加您原来的测试代码以生成额外的幻灯片进行比较或测试
    # --- 测试 StyledQuoteBox ---
    slide1_test = prs.slides.add_slide(blank_slide_layout)
    add_styled_quote_box(slide1_test,
                         "想象力比知识更重要，因为知识是有限的，而想象力概括着世界上的一切，推动着进步，并且是知识进化的源泉。",
                         "阿尔伯特·爱因斯坦",
                         style_id="default_quote", top=Inches(1))
    add_styled_quote_box(slide1_test,
                         "保持简单。",
                         "某设计师",
                         style_id="left_accent_bar_quote", top=Inches(3.5), width=Inches(5))

    # --- 测试 KeyDataCallout ---
    slide2_test = prs.slides.add_slide(blank_slide_layout)
    add_key_data_callout(slide2_test, "75%", "用户满意度",
                         style_id="default_data", left=Inches(0.5), top=Inches(1), width=Inches(3), height=Inches(2))
    add_key_data_callout(slide2_test, "3M+", "下载量",
                         style_id="circle_emphasis_data", left=Inches(4), top=Inches(1), width=Inches(2.5),
                         height=Inches(2.5))

    # --- 测试 DecorativeDivider ---
    slide3_test = prs.slides.add_slide(blank_slide_layout)
    title_shape = slide3_test.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_shape.text_frame.text = "第一部分标题 (测试页)"
    add_decorative_divider(slide3_test, top=Inches(1.6), style_id="default_line")
    content_shape = slide3_test.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
    content_shape.text_frame.text = "这是第一部分的一些内容... (测试页)"
    add_decorative_divider(slide3_test, top=Inches(3.5), style_id="dots_accent_line", width=Inches(5), left=Inches(2.5))

    # 保存演示文稿。
    output_filename = "generated_slide_from_data_v3.pptx"
    prs.save(output_filename)
    print(f"演示文稿已保存为 {output_filename}")
